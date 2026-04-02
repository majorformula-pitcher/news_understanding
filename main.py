import os
import io
from fastapi import FastAPI, Request
from fastapi.responses import HTMLResponse, StreamingResponse, JSONResponse
import httpx
import asyncio
import re
import xml.etree.ElementTree as ET
from bs4 import BeautifulSoup
from jinja2 import Template
import anthropic
from supabase import create_client
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import xlsxwriter

app = FastAPI()

# RSS 피드 목록
RSS_FEEDS = [
    {"name": "로봇신문-AI", "url": "https://www.irobotnews.com/rss/S1N2.xml"},       # 0
    {"name": "로봇신문-로봇", "url": "https://www.irobotnews.com/rss/S1N1.xml"},      # 1
    {"name": "전자신문-AI", "url": "http://rss.etnews.com/04046.xml"},                 # 2
    {"name": "전자신문-전자", "url": "http://rss.etnews.com/06061.xml"},               # 3
    {"name": "The AI", "url": "https://www.newstheai.com/rss/allArticle.xml"},         # 4
    {"name": "디지털투데이", "url": "https://www.digitaltoday.co.kr/rss/allArticle.xml"}, # 5
    {"name": "한국경제-IT", "url": "https://www.hankyung.com/feed/it"},                # 6
    {"name": "ZDNet Korea", "url": "https://zdnet.co.kr/feed"},                        # 7
    {"name": "TechCrunch", "url": "https://techcrunch.com/category/artificial-intelligence/feed/"}, # 8
    {"name": "The Verge", "url": "https://www.theverge.com/rss/index.xml"},            # 9
    {"name": "Wired", "url": "https://www.wired.com/feed/category/business/latest/rss"}, # 10
    {"name": "OpenAI", "url": "https://openai.com/news/rss.xml"},                      # 11
    {"name": "AI Jobs", "url": "https://aijobs.net/feed/"},                            # 12
    {"name": "AI (arxiv)", "url": "http://export.arxiv.org/rss/cs.AI"},                # 13
    {"name": "Techmeme", "url": "https://www.techmeme.com/feed.xml"},                  # 14
    {"name": "Hugging Face", "url": "https://huggingface.co/blog/feed.xml"},           # 15
]

# 활성화된 피드 인덱스
ACTIVE_FEED_INDICES = {0, 1, 2, 3, 4, 5, 6, 7, 8, 10, 14, 15}

# Configure Claude API using environment variable
ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY")

if ANTHROPIC_API_KEY:
    client = anthropic.AsyncAnthropic(api_key=ANTHROPIC_API_KEY)
else:
    print("Warning: ANTHROPIC_API_KEY environment variable is not set.")
    client = None

# Configure Supabase
SUPABASE_URL = os.environ.get("SUPABASE_URL")
SUPABASE_KEY = os.environ.get("SUPABASE_KEY")

supabase = None
supabase_error = None
if SUPABASE_URL and SUPABASE_KEY:
    try:
        supabase = create_client(SUPABASE_URL, SUPABASE_KEY)
    except Exception as e:
        supabase_error = f"Supabase 연결 실패: {e}"
        print(f"Warning: {supabase_error}")
else:
    supabase_error = "SUPABASE_URL 또는 SUPABASE_KEY 환경변수가 설정되지 않았습니다."
    print(f"Warning: {supabase_error}")


def save_articles_to_db(articles, publisher=""):
    """뉴스 기사를 Supabase에 저장. publisher와 published_at 포함"""
    global supabase_error
    supabase_error = None
    if not supabase:
        return
    for article in articles:
        try:
            # 제목 끝에 붙은 " - 출처명" 제거 (RSS 피드에서 자동 추가되는 출처)
            clean_title = article["title"]
            for feed in RSS_FEEDS:
                # 피드 이름에서 "-카테고리" 부분 제거하여 출처명 추출 (e.g. "로봇신문-AI" → "로봇신문")
                source_name = feed["name"].split("-")[0].strip()
                suffix = f" - {source_name}"
                if clean_title.endswith(suffix):
                    clean_title = clean_title[:-len(suffix)].strip()
                    break
            row = {
                "title": clean_title,
                "content": article["body"],
                "summary": article.get("summary", ""),
                "url": article["link"],
                "publisher": publisher or article.get("publisher", ""),
            }
            if article.get("pub_date"):
                row["published_at"] = article["pub_date"]
            supabase.table("news-understanding").upsert(
                row,
                on_conflict="url",
            ).execute()
        except Exception as e:
            supabase_error = f"DB 저장 오류: {e}"
            print(supabase_error)


def load_all_articles():
    """Supabase에서 전체 기사 목록 조회"""
    if not supabase:
        return []
    try:
        all_rows = []
        offset = 0
        batch_size = 1000
        while True:
            result = (
                supabase.table("news-understanding")
                .select("*")
                .order("created_at", desc=True)
                .range(offset, offset + batch_size - 1)
                .execute()
            )
            if not result.data:
                break
            all_rows.extend(result.data)
            if len(result.data) < batch_size:
                break
            offset += batch_size
        return all_rows
    except Exception as e:
        print(f"DB 전체 조회 오류: {e}")
        return []

def load_articles_by_publisher(publisher, order_by="published_at"):
    """Supabase에서 특정 제공자의 기사 목록 조회"""
    global supabase_error
    supabase_error = None
    if not supabase:
        return []
    try:
        all_rows = []
        offset = 0
        batch_size = 1000
        while True:
            result = (
                supabase.table("news-understanding")
                .select("*")
                .eq("publisher", publisher)
                .order(order_by, desc=True)
                .range(offset, offset + batch_size - 1)
                .execute()
            )
            if not result.data:
                break
            all_rows.extend(result.data)
            if len(result.data) < batch_size:
                break
            offset += batch_size
        return [
            {
                "title": row["title"],
                "title_eng": row.get("title_eng", ""),
                "body": row["content"],
                "summary": row.get("summary", ""),
                "summary_eng": row.get("summary_eng", ""),
                "link": row["url"],
                "pub_date": row.get("published_at", ""),
                "publisher": row.get("publisher", ""),
                "is_daily": row.get("is_daily", False),
            }
            for row in all_rows
        ]
    except Exception as e:
        supabase_error = f"DB 조회 오류: {e}"
        print(supabase_error)
        return []


def save_custom_articles_to_db(articles, publisher=""):
    """직접 입력-URL/본문 기사를 Supabase에 저장 (RSS 로직 없이 단순 저장)"""
    global supabase_error
    supabase_error = None
    if not supabase:
        return 0
    saved = 0
    for article in articles:
        try:
            row = {
                "title": article.get("title", ""),
                "content": article.get("body", ""),
                "summary": article.get("summary", ""),
                "url": article.get("link", ""),
                "publisher": publisher,
            }
            # pub_date가 있으면 published_at에 저장, 실패해도 무시
            pub_date = article.get("pub_date", "")
            if pub_date:
                row["published_at"] = pub_date
            supabase.table("news-understanding").upsert(
                row,
                on_conflict="url",
            ).execute()
            saved += 1
        except Exception as e:
            # pub_date 형식 문제일 수 있으므로 published_at 없이 재시도
            try:
                row.pop("published_at", None)
                supabase.table("news-understanding").upsert(
                    row,
                    on_conflict="url",
                ).execute()
                saved += 1
            except Exception as e2:
                supabase_error = f"DB 저장 오류: {e2}"
                print(supabase_error)
    return saved


def update_article_summary(url, summary, summary_eng="", title_eng="", title_ko=""):
    """기사 URL로 summary 컬럼 업데이트"""
    if not supabase:
        return False
    try:
        data = {"summary": summary}
        if summary_eng:
            data["summary_eng"] = summary_eng
        if title_eng:
            data["title_eng"] = title_eng
        if title_ko:
            data["title"] = title_ko
        supabase.table("news-understanding").update(
            data
        ).eq("url", url).execute()
        return True
    except Exception as e:
        print(f"Summary 업데이트 오류: {e}")
        return False


def update_article_daily(url, is_daily):
    """기사 URL로 is_daily 컬럼 업데이트"""
    if not supabase:
        return False
    try:
        supabase.table("news-understanding").update(
            {"is_daily": is_daily}
        ).eq("url", url).execute()
        return True
    except Exception as e:
        print(f"is_daily 업데이트 오류: {e}")
        return False


def load_daily_articles():
    """is_daily=True인 기사 목록 조회"""
    if not supabase:
        return []
    try:
        result = (
            supabase.table("news-understanding")
            .select("*")
            .eq("is_daily", True)
            .order("created_at", desc=True)
            .execute()
        )
        return [
            {
                "title": row["title"],
                "title_eng": row.get("title_eng", ""),
                "body": row["content"],
                "summary": row.get("summary", ""),
                "summary_eng": row.get("summary_eng", ""),
                "link": row["url"],
                "pub_date": row.get("published_at", ""),
                "publisher": row.get("publisher", ""),
                "is_daily": True,
            }
            for row in result.data
        ]
    except Exception as e:
        print(f"Daily 기사 조회 오류: {e}")
        return []


def get_news_stats():
    """DB에 저장된 뉴스 통계 조회 (제공자별 기사 수)"""
    if not supabase:
        return {}
    try:
        result = (
            supabase.table("news-understanding")
            .select("publisher")
            .execute()
        )
        stats = {}
        for row in result.data:
            pub = row.get("publisher", "기타")
            stats[pub] = stats.get(pub, 0) + 1
        return stats
    except Exception as e:
        print(f"뉴스 통계 조회 오류: {e}")
        return {}

def _is_english_text(text):
    """텍스트가 영문인지 판별 (한글이 없으면 영문으로 간주)"""
    return not bool(re.search(r'[가-힣]', text))

async def summarize_article(title, body):
    """한국어 요약 반환. 영문 기사면 제목도 한글 번역하여 (summary, title_ko) 튜플 반환."""
    if not ANTHROPIC_API_KEY or not client:
        return "⚠️ API 키가 설정되지 않았습니다. Render.com 설정에서 ANTHROPIC_API_KEY를 추가해주세요.", ""

    if not body or body == "Content not found" or len(body) < 100:
        return "요약할 충분한 본문 내용이 없습니다.", ""

    is_english = _is_english_text(title)

    if is_english:
        prompt = f"""다음 영문 뉴스 기사를 읽고 아래 형식에 정확히 맞춰 한국어로 요약해 주세요.

형식:
제목: <기사 제목을 한국어로 번역>
. 핵심 요약 첫 번째 줄 (한국어, 2줄 이내)
. 핵심 요약 두 번째 줄 (한국어, 2줄 이내)

주의사항:
- 반드시 '제목:'으로 시작하는 한국어 번역 제목 1줄 + '.'으로 시작하는 2개의 한국어 요약 문장을 작성하세요.
- 불필요한 설명 없이 핵심만 전달하세요.
- 마크다운 문법(**, ##, *, # 등)을 절대 사용하지 마세요. 순수 텍스트로만 작성하세요.

기사 제목: {title}
기사 본문: {body}"""
    else:
        prompt = f"""다음 뉴스 기사를 읽고 아래 형식에 정확히 맞춰 요약해 주세요.

형식:
. 핵심 요약 첫 번째 줄 (2줄 이내)
. 핵심 요약 두 번째 줄 (2줄 이내)

주의사항:
- 기사 제목을 포함하지 마세요. 요약 내용만 작성하세요.
- 반드시 '.'으로 시작하는 2개의 문장으로 요약하세요.
- 불필요한 설명 없이 핵심만 전달하세요.
- 마크다운 문법(**, ##, *, # 등)을 절대 사용하지 마세요. 순수 텍스트로만 작성하세요.

기사 제목: {title}
기사 본문: {body}"""

    try:
        response = await client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=1024,
            system="당신은 뉴스 요약 전문가입니다.",
            messages=[{"role": "user", "content": prompt}]
        )
        text = response.content[0].text.strip()
        text = re.sub(r'\*+', '', text)
        text = re.sub(r'#+\s*', '', text)
        text = re.sub(r'`+', '', text)
        text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
        text = text.strip()

        if is_english:
            title_ko = ""
            summary_lines = []
            for line in text.split('\n'):
                stripped = line.strip()
                if stripped.startswith('제목:'):
                    title_ko = stripped[3:].strip()
                elif stripped:
                    summary_lines.append(stripped)
            return '\n'.join(summary_lines), title_ko

        return text, ""
    except Exception as e:
        return f"요약 중 오류가 발생했습니다: {str(e)}", ""

async def summarize_article_eng(title, body):
    """영문 요약 + 영문 제목 반환 (원본이 영문이면 제목은 그대로 사용)"""
    if not ANTHROPIC_API_KEY or not client:
        return "", ""

    if not body or body == "Content not found" or len(body) < 100:
        return "", "", ""

    is_english = _is_english_text(title)

    if is_english:
        # 영문 원본: 제목 한글 번역 + 영문 요약
        prompt = f"""Read the following news article and provide:
1. A Korean translation of the article title (one line)
2. A summary in English (exactly 2 lines, each starting with '. ')

Your output must be EXACTLY 3 lines in this format:
TITLE_KO: <Korean title here>
. <first key point>
. <second key point>

Do NOT include any other text, headers, or markdown syntax.

Article title: {title}
Article body: {body}"""
    else:
        # 한글 원본: 제목 번역 + 요약
        prompt = f"""Read the following news article and provide:
1. An English translation of the article title (one line)
2. A summary in English (exactly 2 lines, each starting with '. ')

Your output must be EXACTLY 3 lines in this format:
TITLE: <English title here>
. <first key point>
. <second key point>

Do NOT include any other text, headers, or markdown syntax.

Article title: {title}
Article body: {body}"""

    try:
        response = await client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=1024,
            system="You are a news summarization and translation expert.",
            messages=[{"role": "user", "content": prompt}]
        )
        text = response.content[0].text.strip()
        text = re.sub(r'\*+', '', text)
        text = re.sub(r'#+\s*', '', text)
        text = re.sub(r'`+', '', text)
        text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
        text = text.strip()
        # 요약 줄 분리: 문장 끝(. ) 뒤에 바로 '. '로 시작하는 새 줄이 붙어있으면 줄바꿈 추가
        text = re.sub(r'(?<=\.)\s+(?=\.\s)', '\n', text)
        text = text.strip()

        if is_english:
            # 영문 원본: TITLE_KO 줄 분리
            title_ko = ""
            summary_lines = []
            for line in text.split('\n'):
                stripped = line.strip()
                if stripped.upper().startswith('TITLE_KO:'):
                    title_ko = stripped[9:].strip()
                elif stripped:
                    summary_lines.append(stripped)
            summary_eng = '\n'.join(summary_lines)
            return summary_eng, title, title_ko  # (영문요약, 영문제목, 한글제목)

        # 한글 원본: TITLE: 줄 분리
        title_eng = ""
        summary_lines = []
        for line in text.split('\n'):
            stripped = line.strip()
            if stripped.upper().startswith('TITLE:'):
                title_eng = stripped[6:].strip()
            elif stripped:
                summary_lines.append(stripped)
        summary_eng = '\n'.join(summary_lines)
        return summary_eng, title_eng, ""
    except Exception as e:
        print(f"영문 요약 오류: {e}")
        return "", "", ""


async def get_news_content(url):
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/131.0.0.0 Safari/537.36",
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
        "Accept-Language": "en-US,en;q=0.9,ko-KR;q=0.8,ko;q=0.7",
        "Accept-Encoding": "gzip, deflate",
        "Cache-Control": "no-cache",
    }
    error_reason = ""
    try:
        async with httpx.AsyncClient(follow_redirects=True, timeout=15.0) as h_client:
            response = await h_client.get(url, headers=headers)
            status_code = response.status_code
            # 메모리 절약: HTML 크기 제한 (1MB)
            html_text = response.text[:1_000_000] if len(response.text) > 1_000_000 else response.text
            if status_code >= 400:
                # 403/401이라도 HTML에 meta 태그가 있을 수 있으므로 파싱 시도
                soup_err = BeautifulSoup(html_text, 'html.parser')
                og_t = soup_err.find('meta', property='og:title')
                og_d = soup_err.find('meta', property='og:description')
                err_title = og_t['content'].strip() if og_t and og_t.get('content') else ""
                err_desc = og_d['content'].strip() if og_d and og_d.get('content') else ""
                if err_title and err_desc:
                    reason = "[페이월/접근 제한] 전체 본문을 가져올 수 없어 요약 정보만 표시합니다."
                    og_img = soup_err.find('meta', property='og:image')
                    err_image = og_img['content'].strip() if og_img and og_img.get('content') else ""
                    return err_title, reason + "\n\n" + err_desc, "", err_image
                if status_code == 403:
                    return "추출 실패", "[추출 실패] HTTP 403 Forbidden — 이 사이트는 봇 접근을 차단하고 있습니다. (페이월 또는 봇 방지)", "", ""
                elif status_code == 401:
                    return "추출 실패", "[추출 실패] HTTP 401 Unauthorized — 로그인이 필요한 페이지입니다.", "", ""
                else:
                    return "추출 실패", f"[추출 실패] HTTP {status_code} — 서버에서 요청을 거부했습니다.", "", ""

        # JavaScript 렌더링 전용 페이지 감지
        if len(html_text.strip()) < 500 and ('javascript' in html_text.lower() or 'noscript' in html_text.lower()):
            error_reason = "[추출 실패] 이 페이지는 JavaScript로 렌더링되어 서버에서 직접 추출할 수 없습니다."
            return "추출 실패", error_reason, "", ""

        soup = BeautifulSoup(html_text, 'html.parser')

        # --- 제목 추출 (meta 태그 우선) ---
        title = ""
        og_title = soup.find('meta', property='og:title')
        if og_title and og_title.get('content', '').strip():
            title = og_title['content'].strip()
        else:
            twitter_title = soup.find('meta', attrs={'name': 'twitter:title'})
            if twitter_title and twitter_title.get('content', '').strip():
                title = twitter_title['content'].strip()

        if not title:
            title_selectors = ['h2.title', 'h2.title_news', 'div.article_title h2', 'h1.article-title', 'h1.headline', 'h1[data-testid]', 'h1', 'h2']
            for selector in title_selectors:
                el = soup.select_one(selector)
                if el and el.get_text(strip=True):
                    title = el.get_text(strip=True)
                    break

        if not title:
            title_tag = soup.find('title')
            title = title_tag.get_text(strip=True) if title_tag else "제목을 찾을 수 없음"

        # --- 날짜 추출 (meta 태그) ---
        pub_date = ""
        for meta_prop in ['article:published_time', 'og:article:published_time', 'datePublished']:
            meta_el = soup.find('meta', property=meta_prop) or soup.find('meta', attrs={'name': meta_prop})
            if meta_el and meta_el.get('content', '').strip():
                pub_date = meta_el['content'].strip()
                break
        if not pub_date:
            time_el = soup.find('time', attrs={'datetime': True})
            if time_el:
                pub_date = time_el['datetime'].strip()

        # --- 본문 추출 ---
        body_selectors = [
            # 한국 뉴스 사이트
            'div.article_txt', 'div.article_body', 'div#articleBody',
            'div#article-view-content-div', 'div.news_cnt_detail_wrap',
            # 국제 뉴스 사이트
            'div.wp-block-post-content', '[itemprop="articleBody"]', 'div.article-body', 'div.article__body',
            'div.story-body', 'div.article-content', 'div.post-content',
            'div.body-content', 'section.article-body',
            'div[data-component="text-block"]',
            # 대학/기관 사이트
            '#center', 'div.field-item', 'div.node-content',
            # 블로그/일반
            'div.entry-content', 'div.blog-content', 'main article',
            'article', 'div.content', 'main',
        ]
        body_element = None
        for selector in body_selectors:
            el = soup.select_one(selector)
            if el:
                # script, style, nav, footer, aside 태그 제거
                for tag in el.find_all(['script', 'style', 'nav', 'footer', 'aside', 'iframe', 'header']):
                    tag.decompose()
                # Hugging Face: 메타 정보만 선택적 제거 (제목, 작성자, 날짜, 카테고리)
                if 'blog-content' in ' '.join(el.get('class', [])):
                    for tag in el.find_all('div', class_='not-prose'):
                        tag.decompose()
                    for tag in el.find_all('div', class_='mb-4'):
                        tag.decompose()
                    h1 = el.find('h1')
                    if h1:
                        # h1 바로 다음 div (날짜/카테고리 정보) 제거
                        next_div = h1.find_next_sibling('div')
                        if next_div and len(next_div.get_text(strip=True)) < 100:
                            next_div.decompose()
                        h1.decompose()
                if len(el.get_text(strip=True)) > 100:
                    body_element = el
                    break

        if body_element:
            body = body_element.get_text(separator='\n', strip=True)
        else:
            # <p> 태그에서 본문 수집 (셀렉터 매칭 실패 시)
            body = ""
            paragraphs = soup.find_all('p')
            p_texts = [p.get_text(strip=True) for p in paragraphs if len(p.get_text(strip=True)) > 40]
            if p_texts:
                body = '\n'.join(p_texts)

            # meta description 폴백 (p 태그도 부족한 경우 보충)
            if len(body) < 100:
                meta_body = ""
                og_desc = soup.find('meta', property='og:description')
                if og_desc and og_desc.get('content', '').strip():
                    meta_body = og_desc['content'].strip()
                else:
                    meta_desc = soup.find('meta', attrs={'name': 'description'})
                    if meta_desc and meta_desc.get('content', '').strip():
                        meta_body = meta_desc['content'].strip()
                if meta_body:
                    body = meta_body + ("\n\n" + body if body else "")

            if not body:
                error_reason = "[추출 실패] 이 페이지에서 뉴스 본문 영역을 찾을 수 없습니다. 페이월, JavaScript 렌더링, 또는 비표준 HTML 구조일 수 있습니다."
                return title or "추출 실패", error_reason, pub_date, ""

        # 텍스트 정제
        body = re.sub(r'Back to Articles\s*', '', body)
        body = re.sub(r'(?:이메일|email|e-mail)\s*[:\s]*\s*[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}', '', body, flags=re.IGNORECASE)
        body = re.sub(r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}', '', body)
        body = re.sub(r'[가-힣]{2,4}\s*기자(?!\w)', '', body)
        body = re.sub(r'\[\s*\]|\(\s*\)', '', body)
        body = '\n'.join([line.strip() for line in body.split('\n') if line.strip()])
        # 메모리 절약: 본문 길이 제한 (5000자)
        if len(body) > 5000:
            body = body[:5000]

        # --- og:image 추출 ---
        image_url = ""
        og_image = soup.find('meta', property='og:image')
        if og_image and og_image.get('content', '').strip():
            image_url = og_image['content'].strip()
        else:
            twitter_image = soup.find('meta', attrs={'name': 'twitter:image'})
            if twitter_image and twitter_image.get('content', '').strip():
                image_url = twitter_image['content'].strip()

        if len(body) < 50:
            # meta description이라도 있으면 그것을 보여줌
            og_desc = soup.find('meta', property='og:description')
            meta_desc_tag = soup.find('meta', attrs={'name': 'description'})
            fallback = ""
            if og_desc and og_desc.get('content', '').strip():
                fallback = og_desc['content'].strip()
            elif meta_desc_tag and meta_desc_tag.get('content', '').strip():
                fallback = meta_desc_tag['content'].strip()
            if fallback:
                return title or "추출 실패", "[페이월/접근 제한] 전체 본문을 가져올 수 없어 요약 정보만 표시합니다.\n\n" + fallback, pub_date, image_url
            error_reason = f"[추출 실패] 본문이 너무 짧습니다 ({len(body)}자). 페이월이 있거나 JavaScript로 렌더링되는 페이지일 수 있습니다."
            return title or "추출 실패", error_reason, pub_date, image_url

        return title, body, pub_date, image_url
    except httpx.TimeoutException:
        return "추출 실패", "[추출 실패] 요청 시간이 초과되었습니다 (15초). 서버가 응답하지 않거나 봇 접근을 차단하고 있을 수 있습니다.", "", ""
    except httpx.ConnectError:
        return "추출 실패", f"[추출 실패] 서버에 연결할 수 없습니다. URL을 확인해주세요: {url}", "", ""
    except Exception as e:
        return "추출 실패", f"[추출 실패] {type(e).__name__}: {e}", "", ""

async def parse_rss_and_fetch_news(rss_url):
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
    }
    articles = []
    try:
        async with httpx.AsyncClient(follow_redirects=True) as h_client:
            response = await h_client.get(rss_url, headers=headers)
            response.raise_for_status()
            content = response.text
            
        root = ET.fromstring(content)
        items = []
        for item in root.findall('.//item'):
            link_elem = item.find('link')
            if link_elem is not None and link_elem.text:
                title_elem = item.find('title')
                desc_elem = item.find('description')
                rss_title = title_elem.text.strip() if title_elem is not None and title_elem.text else ""
                rss_desc = desc_elem.text.strip() if desc_elem is not None and desc_elem.text else ""
                pub_date_elem = item.find('pubDate')
                rss_pub_date = pub_date_elem.text.strip() if pub_date_elem is not None and pub_date_elem.text else ""
                # Techmeme 등 aggregator: description에서 원본 기사 URL 추출
                original_url = None
                if rss_desc and 'techmeme.com' in (link_elem.text or ''):
                    from bs4 import BeautifulSoup as _BS
                    desc_soup = _BS(rss_desc, 'html.parser')
                    # 본문 큰 링크(SPAN > B > A)에서 원본 URL 추출
                    span = desc_soup.find('span')
                    if span:
                        a_tag = span.find('a', href=True)
                        if a_tag and 'techmeme.com' not in a_tag['href']:
                            original_url = a_tag['href']
                    # description에서 텍스트만 추출하여 요약으로 사용
                    rss_desc = desc_soup.get_text(separator=' ', strip=True)
                items.append({
                    'link': original_url or link_elem.text,
                    'rss_title': rss_title,
                    'rss_desc': rss_desc,
                    'rss_pub_date': rss_pub_date,
                    'original_url': original_url,
                })

        # RSS 아이템이 너무 많으면 최근 50개만 처리
        if len(items) > 50:
            items = items[:50]

        # DB에 이미 저장된 URL 목록 조회하여 새 기사만 필터링
        existing_urls = set()
        if supabase:
            try:
                urls = [it['link'] for it in items]
                result = supabase.table("news-understanding").select("url").in_("url", urls).execute()
                existing_urls = {row["url"] for row in result.data}
            except Exception:
                pass
        new_items = [it for it in items if it['link'] not in existing_urls]

        # 새 기사만 웹페이지 스크래핑 (3개씩 배치 처리)
        fetched_results = []
        batch_size = 3
        for b in range(0, len(new_items), batch_size):
            batch = new_items[b:b+batch_size]
            batch_tasks = [get_news_content(it['link']) for it in batch]
            batch_results = await asyncio.gather(*batch_tasks)
            fetched_results.extend(batch_results)

        for (title, body, page_date, image_url), it in zip(fetched_results, new_items):
            # 페이지 접근 실패 시 RSS 데이터로 대체
            if title.startswith("오류 발생:") or not body or body == "본문을 찾을 수 없습니다." or body.startswith("[추출 실패]") or body.startswith("[페이월"):
                title = it['rss_title'] or title
                body = it['rss_desc'] or body
            # RSS pubDate 우선, 없으면 페이지에서 추출한 날짜 사용
            pub_date = it.get('rss_pub_date', '') or page_date
            articles.append({
                'title': title,
                'body': body,
                'link': it['link'],
                'summary': '',
                'pub_date': pub_date,
            })
    except Exception as e:
        raise Exception(f"RSS 파싱 중 오류 발생: {e}")
    return articles

HTML_TEMPLATE = """
<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>뉴스 요약기</title>
    <style>
        * { box-sizing: border-box; margin: 0; padding: 0; }
        body { font-family: 'Pretendard', sans-serif; background-color: #f0f2f5; display: flex; min-height: 100vh; }

        /* 왼쪽 사이드바 */
        .sidebar {
            width: 200px;
            min-width: 200px;
            background: #1a2233;
            padding: 20px 0;
            display: flex;
            flex-direction: column;
            position: fixed;
            top: 0;
            left: 0;
            height: 100vh;
            overflow-y: auto;
            z-index: 100;
        }
        .sidebar-title {
            color: #fff;
            font-size: 16px;
            font-weight: bold;
            text-align: center;
            padding: 10px 15px 20px;
            border-bottom: 1px solid #2a3a50;
            margin-bottom: 10px;
        }
        .feed-tab {
            display: block;
            width: 100%;
            padding: 14px 18px;
            background: none;
            border: none;
            color: #a0b0c8;
            font-size: 14px;
            text-align: left;
            cursor: pointer;
            transition: all 0.2s;
            border-left: 3px solid transparent;
            text-decoration: none;
        }
        .feed-tab:hover {
            background: #253045;
            color: #fff;
        }
        .feed-tab.active {
            background: #253045;
            color: #4a9eff;
            border-left-color: #4a9eff;
            font-weight: bold;
        }

        /* 오른쪽 콘텐츠 영역 */
        .content {
            margin-left: 200px;
            flex: 1;
            padding: 25px 30px;
        }
        .content-header {
            display: flex;
            align-items: center;
            justify-content: space-between;
            margin-bottom: 25px;
        }
        .content-header h1 {
            color: #1a73e8;
            font-size: 24px;
        }
        .ppt-btn {
            display: inline-block;
            padding: 12px 24px;
            background-color: #e67e22;
            color: white;
            text-decoration: none;
            border-radius: 8px;
            font-size: 15px;
            font-weight: bold;
            transition: background 0.3s;
        }
        .ppt-btn:hover { background-color: #cf6d17; }

        .article-date { font-size: 12px; color: #888; background: #f5f6f8; padding: 3px 10px; border-radius: 4px; white-space: nowrap; flex-shrink: 0; }
        .result-item { background: white; margin-bottom: 30px; padding: 25px; border-radius: 12px; box-shadow: 0 2px 8px rgba(0,0,0,0.07); }
        .result-item h2 { margin-top: 0; color: #333; font-size: 22px; border-bottom: 2px solid #f0f2f5; padding-bottom: 12px; }
        .result-item h2 a { text-decoration: none; color: inherit; }

        .article-layout { display: flex; gap: 25px; margin-top: 15px; align-items: flex-start; }
        .article-body { flex: 0 0 auto; width: 50%; font-size: 15px; line-height: 1.8; color: #444; max-height: 500px; overflow-y: auto; padding-right: 15px; }

        .summary-section { background-color: #fff9db; border: 2px solid #fab005; border-radius: 12px; padding: 20px; display: none; }
        .summary-section.visible { display: block; }
        .summary-section h3 { margin-top: 0; color: #e67e22; font-size: 18px; display: flex; align-items: center; gap: 8px; border-bottom: 1px solid #ffe066; padding-bottom: 10px; }
        .summary-content { font-size: 16px; color: #2c3e50; font-weight: 600; line-height: 1.6; white-space: pre-wrap; }
        .summary-title { font-size: 16px; font-weight: bold; color: #333; margin-bottom: 6px; }
        .summary-eng-section { background-color: #fff9db; border: 2px solid #fab005; border-radius: 12px; padding: 20px; margin-top: 15px; }
        .summary-eng-section.hidden { display: none; }
        .summary-wrapper { flex: 0 0 800px; position: sticky; top: 20px; }

        .btn-row { display: flex; gap: 10px; margin-top: 20px; align-items: center; }
        .original-btn { display: inline-block; padding: 8px 18px; border: 1px solid #1a73e8; color: #1a73e8; text-decoration: none; border-radius: 6px; font-size: 13px; transition: all 0.3s; }
        .original-btn:hover { background-color: #1a73e8; color: white; }
        .summarize-btn { display: inline-block; padding: 8px 18px; background-color: #fab005; color: #fff; border: none; border-radius: 6px; font-size: 13px; font-weight: bold; cursor: pointer; transition: all 0.3s; }
        .summarize-btn:hover { background-color: #e6a200; }
        .summarize-btn:disabled { background-color: #ccc; cursor: not-allowed; }

        .error-msg { background: #fff5f5; color: #c92a2a; padding: 15px; border-radius: 8px; border: 1px solid #ffc9c9; margin-bottom: 20px; }
        .result-item.extraction-error { border-left: 4px solid #e03131; background: #fff8f8; }
        .result-item.extraction-error h2 a { color: #c92a2a; }
        .result-item.extraction-partial { border-left: 4px solid #f59f00; background: #fffdf5; }
        .error-detail { color: #c92a2a; font-size: 14px; padding: 12px 16px; background: #fff0f0; border-radius: 6px; border: 1px solid #ffc9c9; white-space: pre-wrap; line-height: 1.6; }
        .paywall-notice { color: #e67700; font-size: 13px; padding: 8px 12px; background: #fff9e6; border-radius: 6px; border: 1px solid #ffe066; margin-bottom: 8px; }
        .loading-overlay { display: none; text-align: center; color: #1a73e8; font-weight: bold; font-size: 18px; padding: 40px 0; }
        .collect-spinner {
            width: 48px; height: 48px; margin: 0 auto 20px;
            border: 5px solid #e8f0fe; border-top: 5px solid #1a73e8;
            border-radius: 50%; animation: spin 1s linear infinite;
        }
        .collect-spinner.done { animation: none; border-color: #27ae60; border-top-color: #27ae60; }
        @keyframes spin { to { transform: rotate(360deg); } }
        .btn-spinner { display: inline-block; width: 14px; height: 14px; border: 2px solid rgba(255,255,255,0.4); border-top-color: #fff; border-radius: 50%; animation: spin 0.8s linear infinite; vertical-align: middle; margin-right: 6px; }
        .db-error { font-size: 13px; color: #c92a2a; background: #fff5f5; border: 1px solid #ffc9c9; border-radius: 8px; padding: 8px 15px; margin-bottom: 15px; }
        .empty-state { text-align: center; color: #888; padding: 60px 20px; font-size: 18px; }

        .home-tab { border-bottom: 1px solid #2a3a50; margin-bottom: 5px; font-weight: 600; }
        .home-screen { text-align: center; padding: 80px 20px 40px; }
        .home-icon { font-size: 64px; margin-bottom: 20px; }
        .home-title { font-size: 32px; color: #1a73e8; margin-bottom: 12px; font-weight: 700; }
        .home-desc { font-size: 18px; color: #666; margin-bottom: 50px; }
        .home-feeds { display: flex; flex-wrap: wrap; gap: 14px; justify-content: center; max-width: 800px; margin: 0 auto 40px; }
        .home-feed-card {
            padding: 14px 24px; background: #fff; border: 1px solid #e0e0e0; border-radius: 10px;
            color: #333; text-decoration: none; font-size: 15px; font-weight: 500;
            transition: all 0.2s; box-shadow: 0 1px 4px rgba(0,0,0,0.06);
        }
        .home-feed-card:hover { border-color: #1a73e8; color: #1a73e8; box-shadow: 0 3px 12px rgba(26,115,232,0.15); transform: translateY(-2px); }
        .home-hint { font-size: 14px; color: #aaa; }

        .daily-btn {
            flex-shrink: 0;
            padding: 6px 14px;
            background-color: #27ae60;
            color: #fff;
            border: none;
            border-radius: 6px;
            font-size: 12px;
            font-weight: bold;
            cursor: pointer;
            transition: all 0.3s;
            white-space: nowrap;
        }
        .daily-btn:hover { background-color: #219a52; }
        .daily-btn:disabled { background-color: #95a5a6; cursor: not-allowed; }
        .daily-btn.selected { background-color: #95a5a6; cursor: pointer; }
        .daily-btn.selected:hover { background-color: #7f8c8d; }

        .daily-item {
            background: white; margin-bottom: 20px; padding: 20px 25px; border-radius: 12px;
            box-shadow: 0 2px 8px rgba(0,0,0,0.07); position: relative;
            cursor: grab; transition: opacity 0.2s, transform 0.2s;
        }
        .daily-item:active { cursor: grabbing; }
        .daily-item.dragging { opacity: 0.4; transform: scale(0.98); }
        .daily-item.drag-over-top { border-top: 3px solid #1a73e8; margin-top: -3px; }
        .daily-item.drag-over-bottom { border-bottom: 3px solid #1a73e8; margin-bottom: -3px; }
        .daily-item h3 { color: #333; font-size: 18px; margin-bottom: 8px; padding-right: 80px; }
        .daily-item h3 a { text-decoration: none; color: inherit; }
        .daily-item .daily-summary { font-size: 15px; color: #2c3e50; line-height: 1.7; white-space: pre-wrap; }
        .daily-item .daily-remove {
            position: absolute; top: 15px; right: 15px;
            background: #e74c3c; color: #fff; border: none; border-radius: 50%;
            width: 28px; height: 28px; font-size: 16px; cursor: pointer; line-height: 28px; text-align: center;
        }
        .daily-item .daily-order {
            display: inline-block; background: #1a73e8; color: #fff; border-radius: 50%;
            width: 24px; height: 24px; text-align: center; line-height: 24px; font-size: 13px;
            font-weight: bold; margin-right: 8px;
        }
        .daily-source {
            display: inline-block; background: #e8f0fe; color: #1a73e8; font-size: 12px;
            font-weight: 600; padding: 3px 10px; border-radius: 12px; margin-bottom: 8px;
        }

        @media (max-width: 900px) {
            .sidebar { width: 60px; min-width: 60px; }
            .sidebar-title { font-size: 12px; padding: 10px 5px 15px; }
            .feed-tab { font-size: 11px; padding: 10px 8px; }
            .content { margin-left: 60px; padding: 15px; }
            .article-layout { flex-direction: column; }
            .article-layout { flex-direction: column; }
            .article-body { width: 100%; border-bottom: 1px solid #eee; padding-bottom: 15px; padding-right: 0; }
            .summary-section { position: static; flex: 0 0 100%; }
        }
    </style>
</head>
<body>
    <!-- 왼쪽 사이드바: 뉴스 제공자 탭 -->
    <nav class="sidebar">
        <div class="sidebar-title">뉴스 제공자</div>
        <a href="javascript:void(0)"
           class="feed-tab home-tab {% if active_feed is none %}active{% endif %}">
            Home
        </a>
        {% for feed in feeds %}
        {% if loop.index0 in active_indices %}
        <a href="javascript:void(0)"
           class="feed-tab {% if active_feed == loop.index0 %}active{% endif %}"
           data-feed-idx="{{ loop.index0 }}"
           data-feed-name="{{ feed.name }}"
           onclick="loadFeedFromDB({{ loop.index0 }}, '{{ feed.name }}')">
            {{ feed.name }}
        </a>
        {% endif %}
        {% endfor %}
        <a href="javascript:void(0)"
           class="feed-tab"
           onclick="showCustomInput()"
           data-feed-idx="custom">
            직접 입력-URL
        </a>
        <a href="javascript:void(0)"
           class="feed-tab"
           onclick="showCustomManual()"
           data-feed-idx="custom-manual">
            직접 입력-본문
        </a>
    </nav>

    <!-- 오른쪽 콘텐츠 영역 -->
    <main class="content">
        {% if db_error %}
        <div class="db-error">Supabase 오류: {{ db_error }}</div>
        {% endif %}

        <!-- 피드 콘텐츠 헤더 -->
        <div class="content-header" id="feed-header" style="display:none;">
            <h1 id="feed-header-title"></h1>
        </div>

        <div id="loading" class="loading-overlay">뉴스를 불러오는 중입니다...</div>

        <!-- 피드 기사 목록 (JS로 렌더링) -->
        <div id="feed-articles"></div>

        <!-- 직접 입력 섹션 -->
        <div id="custom-section" style="display:none;">
            <div class="content-header">
                <h1>직접 입력</h1>
            </div>
            <div style="padding:20px 0;">
                <p style="color:#888;font-size:14px;margin-bottom:12px;">뉴스 URL을 한 줄에 하나씩 입력하세요 (최대 20개)</p>
                <textarea id="custom-urls" rows="6" style="width:50%;padding:12px;border:1px solid #ddd;border-radius:8px;font-size:14px;resize:vertical;box-sizing:border-box;" placeholder="https://example.com/news/article1&#10;https://example.com/news/article2">https://themiilk.com/articles/aeaf8dcc6?u=12eaa86b&amp;t=ab2139228&amp;from=</textarea>
                <div style="margin-top:12px;display:flex;gap:10px;">
                    <button class="summarize-btn" onclick="fetchCustomUrls()" id="custom-fetch-btn" style="padding:12px 24px;font-size:15px;">기사 가져오기</button>
                </div>
            </div>
            <div id="custom-articles"></div>
        </div>

        <!-- 직접 입력-본문 섹션 -->
        <div id="custom-manual-section" style="display:none;">
            <div class="content-header">
                <h1>직접 입력-본문</h1>
            </div>
            <div style="padding:20px 0;">
                <p style="color:#888;font-size:14px;margin-bottom:12px;">뉴스 제목과 본문을 직접 입력하세요</p>
                <input id="manual-title" type="text" style="width:50%;padding:12px;border:1px solid #ddd;border-radius:8px;font-size:14px;box-sizing:border-box;margin-bottom:8px;" placeholder="뉴스 제목">
                <textarea id="manual-body" rows="10" style="width:50%;padding:12px;border:1px solid #ddd;border-radius:8px;font-size:14px;resize:vertical;box-sizing:border-box;" placeholder="뉴스 본문 내용을 붙여넣으세요"></textarea>
                <input id="manual-url" type="text" style="width:50%;padding:12px;border:1px solid #ddd;border-radius:8px;font-size:14px;box-sizing:border-box;margin-top:8px;" placeholder="원문 URL">
                <div style="margin-top:12px;">
                    <button class="summarize-btn" onclick="saveManualArticle()" id="manual-fetch-btn" style="padding:12px 24px;font-size:15px;">기사 가져오기</button>
                </div>
            </div>
            <div id="manual-articles"></div>
        </div>

        <!-- Home / Daily News -->
        <div id="home-section">
            <!-- 뉴스 수집 상태 메시지 -->
            <div id="collect-status" style="text-align:center;padding:40px 20px;display:none;background:#f0f4ff;border-radius:12px;margin-bottom:20px;">
                <div id="collect-spinner" class="collect-spinner"></div>
                <h2 id="collect-status-text" style="color:#1a73e8;font-size:22px;margin-bottom:10px;">뉴스 정보를 수집 중입니다...</h2>
                <p id="collect-status-sub" style="color:#888;font-size:15px;white-space:pre-wrap;text-align:center;max-width:800px;margin:0 auto;">잠시만 기다려주세요</p>
            </div>
            <h1 style="color:#1a73e8;font-size:32px;text-align:center;margin-bottom:15px;">Daily News</h1>
            <div class="content-header" style="justify-content:space-between;">
                <div style="display:flex;gap:10px;flex-wrap:wrap;">
                    <button class="summarize-btn" onclick="collectNews()" id="collect-btn" style="padding:12px 24px;font-size:15px;">뉴스 업데이트</button>
                    <button class="ppt-btn" onclick="downloadDailyPPT()" id="daily-ppt-btn" style="display:none;">PPT 다운로드</button>
                    <button onclick="clearDaily()" id="daily-clear-btn" style="display:none;padding:12px 24px;background:#e74c3c;color:#fff;border:none;border-radius:8px;font-size:15px;font-weight:bold;cursor:pointer;">선택된 뉴스 전체 해제</button>
                </div>
                <div style="display:flex;gap:10px;flex-wrap:wrap;">
                    <button onclick="exportExcel()" style="padding:12px 24px;background:#217346;color:#fff;border:none;border-radius:8px;font-size:15px;font-weight:bold;cursor:pointer;">DB 엑셀 Export</button>
                    <button onclick="resetDatabase()" style="padding:12px 24px;background:#95a5a6;color:#fff;border:none;border-radius:8px;font-size:15px;font-weight:bold;cursor:pointer;">DB 초기화</button>
                </div>
            </div>
            <div id="last-update-time" style="display:none;padding:8px 0;font-size:13px;color:#888;text-align:left;"></div>
            <div id="daily-empty" class="home-screen" style="padding-top:40px;">
                <div class="home-icon">📰</div>
                <h2 class="home-title">Daily News가 비어 있습니다</h2>
                <p class="home-desc">왼쪽 메뉴에서 뉴스 제공자를 선택한 후<br>"Daily News 로 선택" 버튼을 눌러 기사를 추가하세요</p>
            </div>
            <div id="daily-list"></div>
        </div>
    </main>

    <script>
    function formatPubDate(raw) {
        if (!raw) return '';
        try {
            // 타임존 없는 날짜(예: 2026-03-30 13:50:29)는 이미 KST이므로 그대로 파싱
            var s = String(raw).trim();
            // YYYY-MM-DD HH:MM:SS 형식 직접 처리
            var match = s.match(/^(\\d{4})-(\\d{2})-(\\d{2})\\s+(\\d{2}):(\\d{2})/);
            if (match) {
                return match[1] + '.' + match[2] + '.' + match[3] + ' ' + match[4] + ':' + match[5];
            }
            // RFC 2822 등 타임존이 포함된 형식은 Date 객체로 KST 변환
            const d = new Date(s);
            if (isNaN(d.getTime())) return raw;
            // KST(UTC+9)로 변환
            const kst = new Date(d.getTime() + 9 * 60 * 60 * 1000);
            const y = kst.getUTCFullYear();
            const m = String(kst.getUTCMonth() + 1).padStart(2, '0');
            const day = String(kst.getUTCDate()).padStart(2, '0');
            const h = String(kst.getUTCHours()).padStart(2, '0');
            const min = String(kst.getUTCMinutes()).padStart(2, '0');
            return y + '.' + m + '.' + day + ' ' + h + ':' + min;
        } catch { return raw; }
    }

    // 현재 로드된 기사 배열 (피드별로 갱신됨)
    let articles = [];
    let currentFeedName = '';
    let newsCollected = true;

    // 이스케이프 함수
    function escapeHtml(str) {
        if (!str) return '';
        return str.replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;').replace(/"/g,'&quot;');
    }

    // Daily News 관리 (DB 기반)
    var dailyNewsCache = [];

    function getDailyNews() {
        return dailyNewsCache;
    }

    async function loadDailyNewsFromDB() {
        try {
            const res = await fetch('/api/daily-articles');
            const data = await res.json();
            dailyNewsCache = (data.articles || []).map(a => ({
                title: a.title, title_eng: a.title_eng || '', link: a.link, summary: a.summary,
                summary_eng: a.summary_eng || '', source: a.publisher
            }));
        } catch (e) { dailyNewsCache = []; }
    }

    // 섹션 표시 관리
    var isCollecting = false;

    function showSection(section) {
        document.getElementById('feed-header').style.display = 'none';
        document.getElementById('feed-articles').innerHTML = '';
        document.getElementById('home-section').style.display = 'none';
        document.getElementById('custom-section').style.display = 'none';
        document.getElementById('custom-manual-section').style.display = 'none';
        document.getElementById('loading').style.display = 'none';

        if (section === 'home' || section === 'collect') {
            document.getElementById('home-section').style.display = '';
        } else if (section === 'feed') {
            document.getElementById('feed-header').style.display = '';
        } else if (section === 'custom') {
            document.getElementById('custom-section').style.display = '';
        } else if (section === 'custom-manual') {
            document.getElementById('custom-manual-section').style.display = '';
        }

        // collect-status는 isCollecting 상태에 따라 표시
        document.getElementById('collect-status').style.display = isCollecting ? '' : 'none';
    }

    // 사이드바 active 상태 업데이트
    function setActiveTab(feedIdx) {
        document.querySelectorAll('.feed-tab').forEach(tab => tab.classList.remove('active'));
        if (feedIdx === null) {
            document.querySelector('.feed-tab.home-tab').classList.add('active');
        } else {
            const tab = document.querySelector('.feed-tab[data-feed-idx="' + feedIdx + '"]');
            if (tab) tab.classList.add('active');
        }
    }

    // 직접 입력-URL 탭 클릭
    async function showCustomInput() {
        setActiveTab('custom');
        currentFeedName = '직접 입력-URL';
        showSection('custom');
        await loadCustomArticlesFromDB('직접 입력-URL', 'custom-articles');
    }

    // 직접 입력-본문 탭 클릭
    async function showCustomManual() {
        setActiveTab('custom-manual');
        currentFeedName = '직접 입력-본문';
        showSection('custom-manual');
        await loadCustomArticlesFromDB('직접 입력-본문', 'manual-articles');
    }

    // DB에서 직접 입력 기사 로드
    async function loadCustomArticlesFromDB(publisher, containerId) {
        var container = document.getElementById(containerId);
        container.innerHTML = '<div class="loading-overlay" style="display:block;">기사를 불러오는 중...</div>';
        try {
            var res = await fetch('/api/articles?publisher=' + encodeURIComponent(publisher) + '&order_by=created_at');
            var data = await res.json();
            articles = data.articles || [];
            renderArticleList(containerId);
        } catch (e) {
            container.innerHTML = '<div class="error-msg">기사를 불러오는 중 오류가 발생했습니다.</div>';
        }
    }

    // 직접 입력-본문: 기사 저장
    async function saveManualArticle() {
        var title = document.getElementById('manual-title').value.trim();
        var body = document.getElementById('manual-body').value.trim();
        var url = document.getElementById('manual-url').value.trim();
        if (!title) { alert('제목을 입력해주세요.'); return; }
        if (!body) { alert('본문을 입력해주세요.'); return; }
        if (!url) { alert('URL을 입력해주세요.'); return; }
        var btn = document.getElementById('manual-fetch-btn');
        btn.disabled = true;
        btn.textContent = '저장 중...';
        try {
            await fetch('/api/save-custom-articles', {
                method: 'POST',
                headers: { 'Content-Type': 'application/json' },
                body: JSON.stringify({ articles: [{ title: title, body: body, link: url }], publisher: '직접 입력-본문' })
            });
            // 입력 필드 초기화
            document.getElementById('manual-title').value = '';
            document.getElementById('manual-body').value = '';
            document.getElementById('manual-url').value = '';
            // DB에서 다시 로드
            await loadCustomArticlesFromDB('직접 입력-본문', 'manual-articles');
        } catch (e) {
            alert('저장 중 오류가 발생했습니다.');
        }
        btn.disabled = false;
        btn.textContent = '기사 가져오기';
    }

    // 직접 입력-URL: URL에서 기사 가져오기
    async function fetchCustomUrls() {
        const textarea = document.getElementById('custom-urls');
        const btn = document.getElementById('custom-fetch-btn');
        const urls = textarea.value.trim().split(String.fromCharCode(10)).map(u => u.trim()).filter(u => u.length > 0);

        if (urls.length === 0) {
            alert('URL을 입력해주세요.');
            return;
        }

        btn.disabled = true;
        btn.textContent = '가져오는 중...';
        document.getElementById('custom-articles').innerHTML = '<div class="loading-overlay" style="display:block;">기사를 가져오는 중입니다...</div>';

        try {
            const res = await fetch('/api/fetch-urls', {
                method: 'POST',
                headers: { 'Content-Type': 'application/json' },
                body: JSON.stringify({ urls: urls })
            });
            const data = await res.json();
            var fetched = (data.articles || []).map(a => ({
                title: a.title,
                body: a.body,
                link: a.link,
                summary: '',
                pub_date: a.pub_date || '',
                publisher: '직접 입력-URL',
                error: a.error,
                paywall: a.paywall
            }));
            // DB에 저장
            var validArticles = fetched.filter(a => !a.error);
            if (validArticles.length > 0) {
                await fetch('/api/save-custom-articles', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ articles: validArticles, publisher: '직접 입력-URL' })
                });
            }
            // DB에서 다시 로드 (created_at 순서로)
            await loadCustomArticlesFromDB('직접 입력-URL', 'custom-articles');
        } catch (e) {
            document.getElementById('custom-articles').innerHTML = '<div class="error-msg">기사를 가져오는 중 오류가 발생했습니다.</div>';
        }
        btn.disabled = false;
        btn.textContent = '기사 가져오기';
    }

    // 기사 목록 렌더링 (직접 입력-URL, 직접 입력-본문 공용)
    function renderArticleList(containerId) {
        const container = document.getElementById(containerId);
        if (articles.length === 0) {
            container.innerHTML = '<div class="empty-state">저장된 기사가 없습니다.</div>';
            return;
        }

        container.innerHTML = articles.map((a, i) => {
            if (a.error) {
                return '<div class="result-item" style="opacity:0.6;">' +
                    '<h2><a href="' + escapeHtml(a.link) + '" target="_blank">' + escapeHtml(a.link) + '</a></h2>' +
                    '<p style="color:#e74c3c;">' + escapeHtml(a.body) + '</p>' +
                '</div>';
            }
            const isSelected = a.is_daily;
            const hasSummary = a.summary && a.summary.trim().length > 0;
            const dateFormatted = formatPubDate(a.pub_date);
            return '<div class="result-item" id="article-' + i + '">' +
                '<h2 style="display:flex;align-items:center;justify-content:space-between;gap:12px;">' +
                    '<span style="flex:1;">' + escapeHtml(a.title) + '</span>' +
                    (dateFormatted ? '<span class="article-date">' + dateFormatted + '</span>' : '') +
                    '<button class="daily-btn' + (isSelected ? ' selected' : '') + '" onclick="selectForDaily(' + i + ')" id="daily-btn-' + i + '">' + (isSelected ? '선택됨' : 'Daily News 로 선택') + '</button>' +
                '</h2>' +
                '<div class="article-layout">' +
                    '<div class="article-body">' + escapeHtml(a.body) + '</div>' +
                    '<div class="summary-wrapper">' +
                        '<div class="summary-section' + (hasSummary ? ' visible' : '') + '" id="summary-' + i + '">' +
                            '<div class="summary-title">' + escapeHtml(a.title) + '</div>' +
                            '<div class="summary-content">' + (hasSummary ? escapeHtml(a.summary) : '') + '</div>' +
                        '</div>' +
                        '<div class="summary-eng-section' + ((a.summary_eng || a.title_eng) ? '' : ' hidden') + '" id="summary-eng-section-' + i + '">' +
                            '<div class="summary-title">' + escapeHtml(a.title_eng || '') + '</div>' +
                            '<div class="summary-content">' + escapeHtml(a.summary_eng || '') + '</div>' +
                        '</div>' +
                    '</div>' +
                '</div>' +
                '<div class="btn-row">' +
                    '<a href="' + escapeHtml(a.link) + '" target="_blank" class="original-btn">원문 보기</a>' +
                    '<button class="summarize-btn" onclick="summarizeArticle(' + i + ')"' + (hasSummary ? ' disabled' : '') + '>' + (hasSummary ? '요약 완료' : '요약하기') + '</button>' +
                '</div>' +
            '</div>';
        }).join('');
    }

    // DB에서 기사 로드 후 렌더링
    async function loadFeedFromDB(feedIdx, feedName) {
        setActiveTab(feedIdx);
        currentFeedName = feedName;
        showSection('feed');
        document.getElementById('feed-header-title').textContent = feedName;
        document.getElementById('loading').style.display = 'block';

        try {
            const res = await fetch('/api/articles?publisher=' + encodeURIComponent(feedName));
            const data = await res.json();
            articles = data.articles || [];
            renderFeedArticles();
        } catch (e) {
            document.getElementById('feed-articles').innerHTML = '<div class="error-msg">기사를 불러오는 중 오류가 발생했습니다.</div>';
        }
        document.getElementById('loading').style.display = 'none';
    }

    // 피드 기사 렌더링
    function renderFeedArticles() {
        const container = document.getElementById('feed-articles');
        if (articles.length === 0) {
            container.innerHTML = '<div class="empty-state">이 제공자의 뉴스가 없습니다.</div>';
            return;
        }

        container.innerHTML = articles.map((a, i) => {
            const isSelected = a.is_daily;
            const hasSummary = a.summary && a.summary.trim().length > 0;
            const dateFormatted = formatPubDate(a.pub_date);
            return '<div class="result-item" id="article-' + i + '">' +
                '<h2 style="display:flex;align-items:center;justify-content:space-between;gap:12px;">' +
                    '<span style="flex:1;">' + escapeHtml(a.title) + '</span>' +
                    (dateFormatted ? '<span class="article-date">' + dateFormatted + '</span>' : '') +
                    '<button class="daily-btn' + (isSelected ? ' selected' : '') + '" onclick="selectForDaily(' + i + ')" id="daily-btn-' + i + '">' + (isSelected ? '선택됨' : 'Daily News 로 선택') + '</button>' +
                '</h2>' +
                '<div class="article-layout">' +
                    '<div class="article-body">' + escapeHtml(a.body) + '</div>' +
                    '<div class="summary-wrapper">' +
                        '<div class="summary-section' + (hasSummary ? ' visible' : '') + '" id="summary-' + i + '">' +
                            '<div class="summary-title">' + escapeHtml(a.title) + '</div>' +
                            '<div class="summary-content">' + (hasSummary ? escapeHtml(a.summary) : '') + '</div>' +
                        '</div>' +
                        '<div class="summary-eng-section' + ((a.summary_eng || a.title_eng) ? '' : ' hidden') + '" id="summary-eng-section-' + i + '">' +
                            '<div class="summary-title">' + escapeHtml(a.title_eng || '') + '</div>' +
                            '<div class="summary-content">' + escapeHtml(a.summary_eng || '') + '</div>' +
                        '</div>' +
                    '</div>' +
                '</div>' +
                '<div class="btn-row">' +
                    '<a href="' + escapeHtml(a.link) + '" target="_blank" class="original-btn">원문 보기</a>' +
                    '<button class="summarize-btn" onclick="summarizeArticle(' + i + ')"' + (hasSummary ? ' disabled' : '') + '>' + (hasSummary ? '요약 완료' : '요약하기') + '</button>' +
                '</div>' +
            '</div>';
        }).join('');
    }

    // Daily News 로 선택 버튼 클릭 (토글)
    async function selectForDaily(idx) {
        const btn = document.getElementById('daily-btn-' + idx);
        const article = articles[idx];
        const summaryDiv = document.getElementById('summary-' + idx);
        const contentDiv = summaryDiv.querySelector('.summary-content');

        // 이미 선택된 경우 → 선택 해제
        if (article.is_daily) {
            btn.disabled = true;
            btn.textContent = '해제 중...';
            try {
                await fetch('/api/toggle-daily', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ link: article.link, is_daily: false })
                });
                articles[idx].is_daily = false;
                btn.textContent = 'Daily News 로 선택';
                btn.classList.remove('selected');
                await loadDailyNewsFromDB();
                renderDailyList();
            } catch (e) {
                btn.textContent = '선택됨';
            }
            btn.disabled = false;
            return;
        }

        // DB에서 로드된 요약 또는 화면에 표시된 요약 확인
        const existingSummary = (article.summary && article.summary.trim().length > 0) ? article.summary :
            (summaryDiv.classList.contains('visible') ? contentDiv.textContent : null);

        if (existingSummary && existingSummary !== 'AI가 요약하는 중입니다...' && existingSummary !== '요약 중 오류가 발생했습니다.') {
            btn.disabled = true;
            // 영문 요약이 없으면 요약 실행
            if (!article.summary_eng) {
                btn.textContent = '영문 요약 중...';
                try {
                    const res = await fetch('/api/summarize', {
                        method: 'POST',
                        headers: { 'Content-Type': 'application/json' },
                        body: JSON.stringify({ title: article.title, body: article.body, link: article.link, publisher: currentFeedName })
                    });
                    const data = await res.json();
                    if (data.summary) { articles[idx].summary = data.summary; var cd = document.querySelector('#summary-' + idx + ' .summary-content'); if (cd) cd.textContent = data.summary; }
                    if (data.summary_eng) articles[idx].summary_eng = data.summary_eng;
                    if (data.title_eng) articles[idx].title_eng = data.title_eng;
                    if (data.title_ko) { articles[idx].title = data.title_ko; var a2 = document.querySelector('#article-' + idx + ' h2 a'); if (a2) a2.textContent = data.title_ko; var st2 = document.querySelector('#summary-' + idx + ' .summary-title'); if (st2) st2.textContent = data.title_ko; }
                    var engSec = document.getElementById('summary-eng-section-' + idx);
                    if (engSec && (data.summary_eng || data.title_eng)) {
                        engSec.classList.remove('hidden');
                        engSec.querySelector('.summary-title').textContent = data.title_eng || '';
                        engSec.querySelector('.summary-content').textContent = data.summary_eng || '';
                    }
                } catch (e) {}
            }
            btn.textContent = '선택 중...';
            try {
                await fetch('/api/toggle-daily', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ link: article.link, is_daily: true })
                });
                articles[idx].is_daily = true;
                btn.textContent = '선택됨';
                btn.classList.add('selected');
                await loadDailyNewsFromDB();
                renderDailyList();
            } catch (e) {
                btn.textContent = 'Daily News 로 선택';
            }
            btn.disabled = false;
            return;
        }

        btn.disabled = true;
        btn.textContent = '요약 중...';

        summaryDiv.classList.add('visible');
        contentDiv.textContent = 'AI가 요약하는 중입니다...';

        try {
            const res = await fetch('/api/summarize', {
                method: 'POST',
                headers: { 'Content-Type': 'application/json' },
                body: JSON.stringify({ title: article.title, body: article.body, link: article.link, publisher: currentFeedName })
            });
            const data = await res.json();
            contentDiv.textContent = data.summary;
            articles[idx].summary = data.summary;
            if (data.summary_eng) articles[idx].summary_eng = data.summary_eng;
            if (data.title_eng) articles[idx].title_eng = data.title_eng;
            if (data.title_ko) { articles[idx].title = data.title_ko; var ab = document.querySelector('#article-' + idx + ' h2 a'); if (ab) ab.textContent = data.title_ko; var stb = document.querySelector('#summary-' + idx + ' .summary-title'); if (stb) stb.textContent = data.title_ko; }
            // 영문 요약 영역 표시
            var engSec2 = document.getElementById('summary-eng-section-' + idx);
            if (engSec2 && (data.summary_eng || data.title_eng)) {
                engSec2.classList.remove('hidden');
                engSec2.querySelector('.summary-title').textContent = data.title_eng || '';
                engSec2.querySelector('.summary-content').textContent = data.summary_eng || '';
            }

            // is_daily = true로 설정
            await fetch('/api/toggle-daily', {
                method: 'POST',
                headers: { 'Content-Type': 'application/json' },
                body: JSON.stringify({ link: article.link, is_daily: true })
            });
            articles[idx].is_daily = true;

            btn.textContent = '선택됨';
            btn.classList.add('selected');
            btn.disabled = false;
            await loadDailyNewsFromDB();
            renderDailyList();

            const sumBtn = document.getElementById('article-' + idx).querySelector('.summarize-btn');
            if (sumBtn) { sumBtn.textContent = '요약 완료'; sumBtn.disabled = true; }
        } catch (e) {
            contentDiv.textContent = '요약 중 오류가 발생했습니다.';
            btn.textContent = 'Daily News 로 선택';
            btn.disabled = false;
        }
    }

    async function summarizeArticle(idx) {
        const articleEl = document.getElementById('article-' + idx);
        const btn = articleEl.querySelector('.summarize-btn');
        const summaryDiv = document.getElementById('summary-' + idx);
        const contentDiv = summaryDiv.querySelector('.summary-content');
        const article = articles[idx];

        btn.disabled = true;
        btn.textContent = '요약 중...';
        summaryDiv.classList.add('visible');
        contentDiv.textContent = 'AI가 요약하는 중입니다...';

        try {
            const res = await fetch('/api/summarize', {
                method: 'POST',
                headers: { 'Content-Type': 'application/json' },
                body: JSON.stringify({ title: article.title, body: article.body, link: article.link, publisher: currentFeedName })
            });
            const data = await res.json();
            contentDiv.textContent = data.summary;
            btn.textContent = '요약 완료';
            // 로컬 articles 배열도 업데이트
            articles[idx].summary = data.summary;
            if (data.summary_eng) articles[idx].summary_eng = data.summary_eng;
            if (data.title_eng) articles[idx].title_eng = data.title_eng;
            if (data.title_ko) { articles[idx].title = data.title_ko; var ac = document.querySelector('#article-' + idx + ' h2 a'); if (ac) ac.textContent = data.title_ko; var stc = document.querySelector('#summary-' + idx + ' .summary-title'); if (stc) stc.textContent = data.title_ko; }
            // 영문 요약 영역 표시
            var engSection = document.getElementById('summary-eng-section-' + idx);
            if (engSection && (data.summary_eng || data.title_eng)) {
                engSection.classList.remove('hidden');
                engSection.querySelector('.summary-title').textContent = data.title_eng || '';
                engSection.querySelector('.summary-content').textContent = data.summary_eng || '';
            }
        } catch (e) {
            contentDiv.textContent = '요약 중 오류가 발생했습니다.';
            btn.textContent = '요약하기';
            btn.disabled = false;
        }
    }

    // Home 화면: DB 뉴스 통계 로드 (비동기, 실패해도 무방)
    function loadNewsStats() {
        var emptyDiv = document.getElementById('daily-empty');
        if (!emptyDiv) return;
        fetch('/api/news-stats')
            .then(function(res) { return res.json(); })
            .then(function(data) {
                if (data.total > 0) {
                    var statLines = '';
                    for (var pub in data.stats) {
                        statLines += '<li>' + escapeHtml(pub) + ': <strong>' + data.stats[pub] + '</strong>건</li>';
                    }
                    emptyDiv.innerHTML =
                        '<div class="home-icon">📰</div>' +
                        '<h2 class="home-title">DB에 총 ' + data.total + '건의 뉴스가 저장되어 있습니다</h2>' +
                        '<ul style="list-style:none;padding:0;margin:12px 0;font-size:15px;">' + statLines + '</ul>' +
                        '<p class="home-desc">왼쪽 메뉴에서 뉴스 제공자를 선택하여 기사를 확인하고<br>"Daily News 로 선택" 버튼을 눌러 기사를 추가하세요</p>';
                } else {
                    emptyDiv.innerHTML =
                        '<div class="home-icon">📰</div>' +
                        '<h2 class="home-title">Daily News가 비어 있습니다</h2>' +
                        '<p class="home-desc">"뉴스 업데이트" 버튼을 눌러 뉴스를 수집한 후<br>왼쪽 메뉴에서 뉴스 제공자를 선택하세요</p>';
                }
            })
            .catch(function() {});
    }

    // Home 화면: 선택된 Daily News 목록 렌더링
    function renderDailyList() {
        var listDiv = document.getElementById('daily-list');
        var emptyDiv = document.getElementById('daily-empty');
        var pptBtn = document.getElementById('daily-ppt-btn');
        var clearBtn = document.getElementById('daily-clear-btn');
        if (!listDiv) return;

        var daily = getDailyNews();
        if (daily.length === 0) {
            if (emptyDiv) emptyDiv.style.display = '';
            if (pptBtn) pptBtn.style.display = 'none';
            if (clearBtn) clearBtn.style.display = 'none';
            listDiv.innerHTML = '';
            loadNewsStats();
            return;
        }
        if (emptyDiv) emptyDiv.style.display = 'none';
        if (pptBtn) pptBtn.style.display = '';
        if (clearBtn) clearBtn.style.display = '';
        listDiv.innerHTML = daily.map(function(item, i) {
            // 요약에서 제목 줄 제거 (. 으로 시작하는 줄만 남김)
            var summaryLines = (item.summary || '').split('\\n').filter(function(line) {
                var trimmed = line.trim();
                return trimmed.startsWith('.');
            });
            var cleanSummary = summaryLines.length > 0 ? summaryLines.join('\\n') : item.summary;
            var engLines = (item.summary_eng || '').split('\\n').filter(function(line) {
                var trimmed = line.trim();
                return trimmed.startsWith('.');
            });
            var cleanSummaryEng = engLines.length > 0 ? engLines.join('\\n') : (item.summary_eng || '');
            return '<div class="daily-item" draggable="true" data-index="' + i + '">' +
                '<button class="daily-remove" onclick="removeDaily(' + i + ')">×</button>' +
                '<h3><span class="daily-order">' + (i + 1) + '</span><a href="' + escapeHtml(item.link) + '" target="_blank">' + escapeHtml(item.title) + '</a></h3>' +
                (item.source ? '<span class="daily-source">' + escapeHtml(item.source) + '</span>' : '') +
                '<div class="daily-summary">' + escapeHtml(cleanSummary) + '</div>' +
                (cleanSummaryEng ? '<div style="margin-top:10px;padding-top:10px;border-top:1px solid #ddd;"><div style="font-weight:bold;margin-bottom:4px;">' + escapeHtml(item.title_eng || '') + '</div><div class="daily-summary">' + escapeHtml(cleanSummaryEng) + '</div></div>' : '') +
            '</div>';
        }).join('');

        // 드래그 앤 드롭 이벤트 바인딩
        var items = listDiv.querySelectorAll('.daily-item');
        var dragSrcIndex = null;

        items.forEach(function(item) {
            item.addEventListener('dragstart', function(e) {
                dragSrcIndex = parseInt(this.dataset.index);
                this.classList.add('dragging');
                e.dataTransfer.effectAllowed = 'move';
            });
            item.addEventListener('dragend', function() {
                this.classList.remove('dragging');
                listDiv.querySelectorAll('.daily-item').forEach(function(el) {
                    el.classList.remove('drag-over-top', 'drag-over-bottom');
                });
            });
            item.addEventListener('dragover', function(e) {
                e.preventDefault();
                e.dataTransfer.dropEffect = 'move';
                var rect = this.getBoundingClientRect();
                var midY = rect.top + rect.height / 2;
                this.classList.remove('drag-over-top', 'drag-over-bottom');
                if (e.clientY < midY) {
                    this.classList.add('drag-over-top');
                } else {
                    this.classList.add('drag-over-bottom');
                }
            });
            item.addEventListener('dragleave', function() {
                this.classList.remove('drag-over-top', 'drag-over-bottom');
            });
            item.addEventListener('drop', function(e) {
                e.preventDefault();
                var targetIndex = parseInt(this.dataset.index);
                if (dragSrcIndex === null || dragSrcIndex === targetIndex) return;
                var rect = this.getBoundingClientRect();
                var midY = rect.top + rect.height / 2;
                var dropIndex = e.clientY < midY ? targetIndex : targetIndex + 1;
                var moved = dailyNewsCache.splice(dragSrcIndex, 1)[0];
                if (dropIndex > dragSrcIndex) dropIndex--;
                dailyNewsCache.splice(dropIndex, 0, moved);
                renderDailyList();
            });
        });
    }

    async function removeDaily(idx) {
        const daily = getDailyNews();
        if (idx < 0 || idx >= daily.length) return;
        const item = daily[idx];
        await fetch('/api/toggle-daily', {
            method: 'POST',
            headers: { 'Content-Type': 'application/json' },
            body: JSON.stringify({ link: item.link, is_daily: false })
        });
        await loadDailyNewsFromDB();
        renderDailyList();
    }

    async function clearDaily() {
        if (!confirm('선택한 Daily News를 모두 해제하시겠습니까?')) return;
        const daily = getDailyNews();
        await Promise.all(daily.map(item =>
            fetch('/api/toggle-daily', {
                method: 'POST',
                headers: { 'Content-Type': 'application/json' },
                body: JSON.stringify({ link: item.link, is_daily: false })
            })
        ));
        await loadDailyNewsFromDB();
        renderDailyList();
    }

    async function exportExcel() {
        if (!confirm('전체 뉴스를 엑셀로 Export 하시겠습니까?')) return;
        try {
            const res = await fetch('/api/export-excel');
            if (!res.ok) throw new Error('엑셀 생성 실패');
            const blob = await res.blob();
            const url = URL.createObjectURL(blob);
            const a = document.createElement('a');
            a.href = url;
            a.download = 'news_export.xlsx';
            a.click();
            URL.revokeObjectURL(url);
        } catch (e) {
            alert('엑셀 다운로드 중 오류가 발생했습니다: ' + e.message);
        }
    }

    async function resetDatabase() {
        if (!confirm('전체 뉴스 DB가 삭제됩니다. 계속하시겠습니까?')) return;
        try {
            const res = await fetch('/api/reset-db', { method: 'POST' });
            const data = await res.json();
            if (res.ok) {
                var msg = 'DB가 초기화되었습니다. 총 ' + data.deleted + '건 삭제됨.';
                if (!data.id_reset) msg += '\\n\\n(ID 초기화 실패: Supabase SQL Editor에서 truncate_news 함수를 생성해주세요)';
                alert(msg);
                loadNewsStats();
            } else {
                alert('DB 초기화 오류: ' + (data.error || '알 수 없는 오류'));
            }
        } catch (e) {
            alert('DB 초기화 중 오류가 발생했습니다: ' + e.message);
        }
    }

    async function downloadDailyPPT() {
        const daily = getDailyNews();
        if (daily.length === 0) { alert('선택된 뉴스가 없습니다.'); return; }
        try {
            const res = await fetch('/api/daily-ppt', {
                method: 'POST',
                headers: { 'Content-Type': 'application/json' },
                body: JSON.stringify({ articles: daily })
            });
            if (!res.ok) throw new Error('PPT 생성 실패');
            const blob = await res.blob();
            const url = URL.createObjectURL(blob);
            const a = document.createElement('a');
            a.href = url;
            var today = new Date(); var ymd = today.getFullYear() + String(today.getMonth()+1).padStart(2,'0') + String(today.getDate()).padStart(2,'0');
            a.download = 'daily_news_' + ymd + '.pptx';
            a.click();
            URL.revokeObjectURL(url);
        } catch (e) {
            alert('PPT 다운로드 중 오류가 발생했습니다.');
        }
    }

    // 직접 입력 탭 비활성화/활성화
    function setCustomTabsDisabled(disabled) {
        var tabs = document.querySelectorAll('.feed-tab[data-feed-idx="custom"], .feed-tab[data-feed-idx="custom-manual"]');
        for (var i = 0; i < tabs.length; i++) {
            tabs[i].style.pointerEvents = disabled ? 'none' : '';
            tabs[i].style.opacity = disabled ? '0.4' : '';
        }
    }

    // 뉴스 업데이트 버튼 클릭
    async function collectNews() {
        const btn = document.getElementById('collect-btn');
        btn.disabled = true;
        setCustomTabsDisabled(true);
        btn.innerHTML = '<span class="btn-spinner"></span>수집 중...';
        isCollecting = true;

        // 업데이트 시작 시간 표시
        var now = new Date();
        var timeStr = now.getFullYear() + '-' + String(now.getMonth()+1).padStart(2,'0') + '-' + String(now.getDate()).padStart(2,'0') + ' ' + String(now.getHours()).padStart(2,'0') + ':' + String(now.getMinutes()).padStart(2,'0') + ':' + String(now.getSeconds()).padStart(2,'0');
        var updateTimeEl = document.getElementById('last-update-time');
        updateTimeEl.textContent = '마지막 뉴스 업데이트: ' + timeStr;
        updateTimeEl.style.display = '';

        showSection('collect');
        var spinner = document.getElementById('collect-spinner');
        spinner.className = 'collect-spinner';
        var statusText = document.getElementById('collect-status-text');
        var statusSub = document.getElementById('collect-status-sub');
        statusText.textContent = '뉴스 정보를 수집 중입니다...';
        statusSub.style.textAlign = 'center';
        statusSub.textContent = '잠시만 기다려주세요';

        try {
            const res = await fetch('/api/collect-news-stream');
            const reader = res.body.getReader();
            const decoder = new TextDecoder();
            var buffer = '';
            var finalData = null;

            while (true) {
                const { done, value } = await reader.read();
                if (done) break;
                buffer += decoder.decode(value, { stream: true });
                var lines = buffer.split('\\n');
                buffer = lines.pop();
                for (var line of lines) {
                    if (!line.startsWith('data: ')) continue;
                    try {
                        var evt = JSON.parse(line.substring(6));
                        if (evt.type === 'progress') {
                            statusSub.textContent = '뉴스 ' + evt.analyzed + '개를 분석했습니다. 뉴스 ' + evt.collected + '개를 DB에 저장했습니다.\\n현재: ' + evt.feed;
                        } else if (evt.type === 'error') {
                            statusSub.textContent += '\\n⚠ ' + evt.feed + ': 오류 발생';
                        } else if (evt.type === 'done') {
                            finalData = evt;
                        }
                    } catch(x) {}
                }
            }

            newsCollected = true;
            spinner.className = 'collect-spinner done';
            if (finalData) {
                statusText.textContent = '뉴스 정보 수집이 완료됐습니다';
                statusSub.textContent = '뉴스 ' + finalData.analyzed + '개를 분석하고, ' + finalData.collected + '개를 DB에 저장했습니다.\\n왼쪽 메뉴에서 뉴스 제공자를 선택하세요.';
                if (finalData.errors && finalData.errors.length > 0) {
                    statusSub.textContent += '\\n\\n⚠ 일부 피드 오류:\\n' + finalData.errors.join('\\n');
                    statusSub.style.textAlign = 'left';
                }
            }
            var delay = (finalData && finalData.errors && finalData.errors.length > 0) ? 6000 : 3000;
            setTimeout(function() {
                isCollecting = false;
                setCustomTabsDisabled(false);
                document.getElementById('collect-status').style.display = 'none';
                renderDailyList();
            }, delay);
        } catch (e) {
            isCollecting = false;
            setCustomTabsDisabled(false);
            spinner.className = 'collect-spinner done';
            statusText.textContent = '뉴스 수집 중 오류가 발생했습니다';
            statusSub.textContent = e.message || '알 수 없는 오류';
        }
        btn.disabled = false;
        btn.textContent = '뉴스 업데이트';
    }

    // Home 탭 클릭
    document.querySelector('.feed-tab.home-tab').addEventListener('click', async function(e) {
        e.preventDefault();
        setActiveTab(null);
        showSection('home');
        await loadDailyNewsFromDB();
        renderDailyList();
    });

    // 초기화: 홈 화면 표시
    showSection('home');
    loadDailyNewsFromDB().then(function() { renderDailyList(); });
    </script>
</body>
</html>
"""

def generate_ppt(articles):
    """뉴스 요약을 PPT로 생성 (슬라이드 1장당 1개 기사: 한글+영문)"""
    prs = Presentation()
    prs.slide_width = Cm(14.001)
    prs.slide_height = Cm(6.002)
    FONT_NAME = 'SamsungOne 700'

    # 타이틀 슬라이드
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    txBox = slide.shapes.add_textbox(Cm(0.5), Cm(1.5), Cm(13.0), Cm(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "뉴스 핵심 요약"
    p.font.name = FONT_NAME
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 115, 232)
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = f"총 {len(articles)}건"
    p2.font.name = FONT_NAME
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(100, 100, 100)
    p2.alignment = PP_ALIGN.CENTER

    def _to_bullet(line):
        """선행 . 을 • 글머리 기호로 변환"""
        if line.startswith('. '):
            return '∙ ' + line[2:]
        if line.startswith('.'):
            return '∙' + line[1:]
        return line

    def _fetch_og_image(page_url):
        """기사 URL에서 og:image를 추출하고 다운로드하여 BytesIO로 반환. 실패 시 None."""
        if not page_url:
            return None
        try:
            with httpx.Client(follow_redirects=True, timeout=10.0) as client:
                # 1) 페이지에서 og:image URL 추출
                resp = client.get(page_url, headers={"User-Agent": "Mozilla/5.0"})
                if resp.status_code >= 400:
                    return None
                soup = BeautifulSoup(resp.text[:200_000], 'html.parser')
                og_img = soup.find('meta', property='og:image')
                if not og_img or not og_img.get('content', '').strip():
                    og_img = soup.find('meta', attrs={'name': 'twitter:image'})
                if not og_img or not og_img.get('content', '').strip():
                    return None
                img_url = og_img['content'].strip()
                # 2) 이미지 다운로드
                img_resp = client.get(img_url, headers={"User-Agent": "Mozilla/5.0"})
                if img_resp.status_code == 200 and img_resp.headers.get("content-type", "").startswith("image"):
                    return io.BytesIO(img_resp.content)
        except Exception:
            pass
        return None

    IMG_SIZE = Cm(2)  # 2cm x 2cm (슬라이드 크기에 맞춤)
    IMG_COL_WIDTH = Cm(2.3)  # 이미지 컬럼 (이미지 + 여백)
    TABLE_LEFT = Cm(0.5)
    TABLE_TOP = Cm(0.2)
    TABLE_WIDTH = Cm(13.0)
    TABLE_HEIGHT = Cm(4.8)
    TITLE_ROW_HEIGHT = Cm(0.8)
    BODY_ROW_HEIGHT = TABLE_HEIGHT - TITLE_ROW_HEIGHT

    def _set_cell_font(cell, text, font_name, font_size, bold=False, color=None):
        """셀에 텍스트 설정 및 폰트 적용"""
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.name = font_name
        p.font.size = font_size
        p.font.bold = bold
        if color:
            p.font.color.rgb = color
        cell.text_frame.word_wrap = True

    def _add_summary_slide(prs, title_text, summary_text, article_url, image_data, is_english=False):
        """요약 슬라이드 1장 생성. 표 레이아웃: Row1=제목, Row2=요약|이미지"""
        from pptx.oxml.ns import qn

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        font_color = RGBColor(0, 0, 0)
        title_color = RGBColor(2, 44, 178)
        has_image = image_data is not None

        # 표 생성: 2행 x 2열
        cols = 2 if has_image else 1
        table_shape = slide.shapes.add_table(2, cols, TABLE_LEFT, TABLE_TOP, TABLE_WIDTH, TABLE_HEIGHT)
        table = table_shape.table

        # 컬럼 폭 설정
        if has_image:
            text_col_width = TABLE_WIDTH - IMG_COL_WIDTH
            table.columns[0].width = int(text_col_width)
            table.columns[1].width = int(IMG_COL_WIDTH)
        else:
            table.columns[0].width = int(TABLE_WIDTH)

        # 행 높이 설정
        table.rows[0].height = int(TITLE_ROW_HEIGHT)
        table.rows[1].height = int(BODY_ROW_HEIGHT)

        # Row 0: 제목 (셀 병합)
        if has_image:
            # 첫 번째 행의 두 셀을 병합
            cell_title_left = table.cell(0, 0)
            cell_title_right = table.cell(0, 1)
            cell_title_left.merge(cell_title_right)
        title_cell = table.cell(0, 0)
        title_cell.text = ""
        p_title = title_cell.text_frame.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = FONT_NAME
        p_title.font.size = Pt(8)
        p_title.font.bold = True
        p_title.font.underline = True
        p_title.font.color.rgb = title_color
        title_cell.text_frame.word_wrap = True
        # 제목 셀 배경색
        title_cell.fill.solid()
        title_cell.fill.fore_color.rgb = RGBColor(230, 230, 230)

        # Row 1, Col 0: 요약 본문
        summary_cell = table.cell(1, 0)
        summary_cell.text = ""
        summary_cell.fill.solid()
        summary_cell.fill.fore_color.rgb = RGBColor(230, 230, 230)
        summary_cell.text_frame.word_wrap = True
        lines = [l.strip() for l in summary_text.split('\n') if l.strip()]
        first = True
        for line in lines:
            if first:
                p = summary_cell.text_frame.paragraphs[0]
                first = False
            else:
                p = summary_cell.text_frame.add_paragraph()
            p.text = _to_bullet(line)
            p.font.name = FONT_NAME
            p.font.size = Pt(7)
            p.font.color.rgb = font_color
            p.line_spacing = Pt(11)
            p.space_before = Pt(2)

        # Row 1, Col 1: 이미지 셀 (이미지를 셀 위에 겹쳐서 배치)
        if has_image:
            img_cell = table.cell(1, 1)
            img_cell.text = ""
            img_cell.fill.solid()
            img_cell.fill.fore_color.rgb = RGBColor(230, 230, 230)
            # 이미지를 표의 오른쪽 셀 영역 중앙에 배치
            image_data.seek(0)
            img_left = TABLE_LEFT + int(table.columns[0].width) + (IMG_COL_WIDTH - IMG_SIZE) // 2
            img_top = TABLE_TOP + TITLE_ROW_HEIGHT + (BODY_ROW_HEIGHT - IMG_SIZE) // 2
            pic = slide.shapes.add_picture(image_data, img_left, img_top, IMG_SIZE, IMG_SIZE)
            # 이미지 모서리 둥글게 (roundRect, adj=16667 = 1/6)
            spPr = pic._element.spPr
            prstGeom = spPr.find(qn('a:prstGeom'))
            if prstGeom is not None:
                spPr.remove(prstGeom)
            prstGeom = spPr.makeelement(qn('a:prstGeom'), {'prst': 'roundRect'})
            avLst = prstGeom.makeelement(qn('a:avLst'), {})
            gd = avLst.makeelement(qn('a:gd'), {'name': 'adj', 'fmla': 'val 16667'})
            avLst.append(gd)
            prstGeom.append(avLst)
            spPr.append(prstGeom)

        # 표 테두리 완전 제거 (noFill)
        tbl = table._tbl
        tblPr = tbl.tblPr
        # tblStyle 제거 (기본 테마 테두리 방지)
        tblStyleId = tblPr.find(qn('a:tblStyleId'))
        if tblStyleId is not None:
            tblPr.remove(tblStyleId)
        # 개별 셀 테두리: noFill로 설정하여 완전히 안 보이게 처리
        for row_idx in range(len(table.rows)):
            for col_idx in range(cols):
                cell = table.cell(row_idx, col_idx)
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                    border = tcPr.find(qn(border_name))
                    if border is not None:
                        tcPr.remove(border)
                    ln = tcPr.makeelement(qn(border_name), {'w': '0'})
                    noFill = ln.makeelement(qn('a:noFill'), {})
                    ln.append(noFill)
                    tcPr.append(ln)

        # URL을 슬라이드 노트에 추가
        if article_url:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = article_url

    # 기사마다 한국어 슬라이드 1장 + 영문 슬라이드 1장
    for article in articles:
        article_url = article.get("link", "")
        image_data = _fetch_og_image(article_url)

        title = article.get("title", "")
        title_eng = article.get("title_eng", "")
        is_eng_article = _is_english_text(title)

        # 한국어 요약 슬라이드: 영문 기사면 title_eng에 한글 번역이 없으므로 title 그대로 사용
        # (DB에서 title이 한글로 업데이트된 경우 그대로 사용)
        ko_title = title
        ko_summary = article.get("summary") or "요약 없음"
        _add_summary_slide(prs, ko_title, ko_summary, article_url, image_data)

        # 영문 요약 슬라이드
        en_title = title_eng if title_eng else title
        en_summary = article.get("summary_eng") or "No English summary"
        _add_summary_slide(prs, en_title, en_summary, article_url, image_data, is_english=True)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


@app.get("/api/collect-news-stream")
async def api_collect_news_stream():
    """SSE로 실시간 진행률을 전송하며 뉴스 수집"""
    import json as _json

    async def event_stream():
        collected = 0
        analyzed = 0
        errors = []
        feed_list = [RSS_FEEDS[idx] for idx in sorted(ACTIVE_FEED_INDICES)]

        for feed in feed_list:
            try:
                articles = await parse_rss_and_fetch_news(feed["url"])
                analyzed += len(articles)
                yield f"data: {_json.dumps({'type': 'progress', 'analyzed': analyzed, 'collected': collected, 'feed': feed['name']})}\n\n"
                if articles:
                    save_articles_to_db(articles, publisher=feed["name"])
                    collected += len(articles)
                    yield f"data: {_json.dumps({'type': 'progress', 'analyzed': analyzed, 'collected': collected, 'feed': feed['name']})}\n\n"
            except Exception as e:
                errors.append(f"{feed['name']}: {e}")
                yield f"data: {_json.dumps({'type': 'error', 'feed': feed['name'], 'message': str(e)})}\n\n"

        yield f"data: {_json.dumps({'type': 'done', 'analyzed': analyzed, 'collected': collected, 'errors': errors})}\n\n"

    return StreamingResponse(event_stream(), media_type="text/event-stream")


@app.get("/api/articles")
async def api_articles(publisher: str = "", order_by: str = "published_at"):
    """DB에서 특정 제공자의 기사 목록 조회"""
    if not publisher:
        return JSONResponse({"articles": []})
    articles = load_articles_by_publisher(publisher, order_by=order_by)
    return JSONResponse({"articles": articles})


@app.post("/api/summarize")
async def api_summarize(request: Request):
    """개별 기사 요약 API (한국어 + 영어 동시 진행)"""
    data = await request.json()
    title = data.get("title", "")
    body = data.get("body", "")
    link = data.get("link", "")
    publisher = data.get("publisher", "")
    # 한국어/영어 요약 동시 실행
    ko_result, eng_result = await asyncio.gather(
        summarize_article(title, body),
        summarize_article_eng(title, body),
    )
    summary, title_ko_from_summary = ko_result
    summary_eng, title_eng, title_ko_from_eng = eng_result
    # 한글 제목: summarize_article에서 번역한 것 우선, 없으면 eng에서 가져옴
    title_ko = title_ko_from_summary or title_ko_from_eng
    # DB 업데이트
    if summary and link:
        updated = update_article_summary(link, summary, summary_eng, title_eng, title_ko)
        if not updated:
            if publisher.startswith("직접 입력"):
                save_custom_articles_to_db([{"title": title, "body": body, "link": link, "summary": summary}], publisher=publisher)
            else:
                save_articles_to_db([{"title": title, "body": body, "link": link, "summary": summary}], publisher=publisher)
            if summary_eng:
                update_article_summary(link, summary, summary_eng, title_eng, title_ko)
    return JSONResponse({"summary": summary, "summary_eng": summary_eng, "title_eng": title_eng, "title_ko": title_ko})


@app.post("/api/toggle-daily")
async def api_toggle_daily(request: Request):
    """기사의 is_daily 토글"""
    data = await request.json()
    link = data.get("link", "")
    is_daily = data.get("is_daily", False)
    if not link:
        return JSONResponse({"error": "link 필요"}, status_code=400)
    ok = update_article_daily(link, is_daily)
    if not ok:
        return JSONResponse({"error": "업데이트 실패"}, status_code=500)
    return JSONResponse({"status": "ok", "is_daily": is_daily})


@app.get("/api/daily-articles")
async def api_daily_articles():
    """is_daily=True인 기사 목록 조회"""
    articles = load_daily_articles()
    return JSONResponse({"articles": articles})


@app.get("/api/news-stats")
async def api_news_stats():
    """DB에 저장된 뉴스 통계 조회"""
    stats = get_news_stats()
    total = sum(stats.values())
    return JSONResponse({"stats": stats, "total": total})


@app.post("/api/fetch-urls")
async def api_fetch_urls(request: Request):
    """URL 목록에서 뉴스 제목/본문 추출"""
    data = await request.json()
    urls = data.get("urls", [])
    fetch_tasks = [get_news_content(url) for url in urls[:20]]
    results = await asyncio.gather(*fetch_tasks)
    articles = []
    for (title, body, pub_date, image_url), url in zip(results, urls[:20]):
        is_error = title == "추출 실패" or body.startswith("[추출 실패]")
        is_paywall = body.startswith("[페이월/접근 제한]")
        articles.append({"title": title, "body": body, "link": url, "error": is_error, "paywall": is_paywall, "pub_date": pub_date})
    return JSONResponse({"articles": articles})


@app.post("/api/save-custom-articles")
async def api_save_custom_articles(request: Request):
    """직접 입력-URL/본문 기사를 DB에 저장"""
    data = await request.json()
    articles = data.get("articles", [])
    publisher = data.get("publisher", "직접 입력-URL")
    saved = save_custom_articles_to_db(articles, publisher=publisher)
    return JSONResponse({"saved": saved})


@app.post("/api/daily-ppt")
async def daily_ppt(request: Request):
    """Daily News 선택 기사들로 PPT 생성 (DB에서 최신 데이터 로드)"""
    articles = load_daily_articles()
    if not articles:
        return JSONResponse({"error": "선택된 기사가 없습니다."}, status_code=400)
    output = generate_ppt(articles)
    from datetime import datetime
    filename = f"daily_news_{datetime.now().strftime('%Y%m%d')}.pptx"
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={filename}"},
    )


@app.post("/api/reset-db")
async def api_reset_db():
    """DB 전체 기사 삭제"""
    if not supabase or not SUPABASE_URL or not SUPABASE_KEY:
        return JSONResponse({"error": "DB 연결이 없습니다."}, status_code=500)
    try:
        # 먼저 건수 확인
        count_result = supabase.table("news-understanding").select("id", count="exact").execute()
        deleted = count_result.count or 0

        # 전체 삭제 (id > 0 조건으로 모든 행 삭제)
        supabase.table("news-understanding").delete().gt("id", 0).execute()

        return JSONResponse({"status": "ok", "deleted": deleted, "id_reset": True})
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/export-excel")
async def api_export_excel():
    """DB 전체 기사를 엑셀로 export"""
    rows = load_all_articles()
    if not rows:
        return JSONResponse({"error": "데이터가 없습니다."}, status_code=400)

    output = io.BytesIO()
    workbook = xlsxwriter.Workbook(output, {'in_memory': True})
    worksheet = workbook.add_worksheet('뉴스')

    # 헤더 스타일
    header_fmt = workbook.add_format({'bold': True, 'bg_color': '#1a73e8', 'font_color': '#ffffff', 'border': 1})
    cell_fmt = workbook.add_format({'text_wrap': True, 'valign': 'top', 'border': 1})

    headers = ['제목', '출처', 'URL', '요약', '본문', '발행일', '수집일']
    for col, h in enumerate(headers):
        worksheet.write(0, col, h, header_fmt)

    worksheet.set_column(0, 0, 40)  # 제목
    worksheet.set_column(1, 1, 15)  # 출처
    worksheet.set_column(2, 2, 30)  # URL
    worksheet.set_column(3, 3, 50)  # 요약
    worksheet.set_column(4, 4, 60)  # 본문
    worksheet.set_column(5, 5, 20)  # 발행일
    worksheet.set_column(6, 6, 20)  # 수집일

    for i, row in enumerate(rows):
        worksheet.write(i + 1, 0, row.get('title', ''), cell_fmt)
        worksheet.write(i + 1, 1, row.get('publisher', ''), cell_fmt)
        worksheet.write(i + 1, 2, row.get('url', ''), cell_fmt)
        worksheet.write(i + 1, 3, row.get('summary', ''), cell_fmt)
        worksheet.write(i + 1, 4, row.get('content', ''), cell_fmt)
        worksheet.write(i + 1, 5, row.get('published_at', ''), cell_fmt)
        worksheet.write(i + 1, 6, row.get('created_at', ''), cell_fmt)

    workbook.close()
    output.seek(0)
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": "attachment; filename=news_export.xlsx"},
    )


@app.get("/", response_class=HTMLResponse)
async def get_index(feed: str = None):
    template = Template(HTML_TEMPLATE)
    active_feed = None

    if feed == "custom":
        active_feed = "custom"
    elif feed is not None:
        try:
            feed_idx = int(feed)
            if 0 <= feed_idx < len(RSS_FEEDS) and feed_idx in ACTIVE_FEED_INDICES:
                active_feed = feed_idx
        except ValueError:
            pass

    default_manual_title = "Anthropic's Claude Code Leak Revealed Unreleased Features"
    default_manual_body = """Anthropic PBC's accidental release of source code for its popular AI coding agent was the result of "process errors" related to the startup's fast product release cycle, according to a senior executive at the company.

Paul Smith, Anthropic's chief commercial officer, said the leak was "absolutely not breaches or hacks," and the mistakes have been addressed. "They're part of the incredibly rapid release cycle that we've had around Claude Code," Smith said in an interview Wednesday.

Anthropic takes the issue "incredibly seriously," he added, noting that the firm has "all the right people focused on addressing it."
The unintended release marked Anthropic's second security slip-up in a matter of days, compromising approximately 1,900 files and 512,000 lines of code related to Claude Code. Last week, Fortune separately reported that Anthropic had been storing thousands of internal files on a publicly accessible system, including a draft blog post that detailed an upcoming model known internally as both "Mythos" and "Capybara."
In a series of posts overnight on X, Claude Code creator Boris Cherny said Anthropic's "deploy process has a few manual steps, and we didn't do one of the steps correctly." He said the company has already "made a few improvements to the automation for next time," with plans for "a couple more on the way."

The code included at least eight unreleased features, according to Abhishek Tiwari, an entrepreneur who quickly coded a website to track the leak using AI. Among them, he listed:

Kairos, a setting that allows Claude to run in the background.
Coordinator Mode, which allows AI to break a task into pieces and delegate them to individual workers.
Auto-Dream, in which Claude reviews what it learns and organizes notes into clean, structured memory files.
Ultraplan, which allows the creation of a separate cloud instance that explores and plans up to 30 minutes at a time.
Anthropic did not comment on the features or its plans. "We're always experimenting with new ideas. 90% don't ship because we don't think they're good enough experiences," Cherny said in a response to a post on X."""
    default_manual_url = "https://www.bloomberg.com/news/articles/2026-04-01/anthropic-scrambles-to-address-leak-of-claude-code-source-code?srnd=phx-technology"

    return HTMLResponse(content=template.render(
        feeds=RSS_FEEDS,
        active_feed=active_feed,
        active_indices=ACTIVE_FEED_INDICES,
        articles=[],
        error=None,
        db_error=supabase_error,
        default_manual_title=default_manual_title,
        default_manual_body=default_manual_body,
        default_manual_url=default_manual_url,
    ))

if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PORT", 8000))
    uvicorn.run("main:app", host="0.0.0.0", port=port)
